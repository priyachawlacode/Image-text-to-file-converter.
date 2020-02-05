from PIL import Image
from fpdf import FPDF
from pptx import Presentation
from pytesseract import image_to_string
import pytesseract
import docx
import sys
import os
import comtypes.client
import numpy as np
import cv2

#CONTRASTING THE IMAGE

# read
img = cv2.imread('test.jpg', cv2.IMREAD_GRAYSCALE)

# increase contrast
pxmin = np.min(img)
pxmax = np.max(img)
imgContrast = (img - pxmin) / (pxmax - pxmin) * 255

# increase line width
kernel = np.ones((3, 3), np.uint8)
imgMorph = cv2.erode(imgContrast, kernel, iterations = 1)

# write
#cv2.imwrite('out.png', imgMorph)



#CONVERTING TO STRING

pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
p=image_to_string( imgMorph,lang='eng')
f=open("text.txt","a")
f.write(p)
f.close()
#print(p)
ch=int(input("enter your choice\n 1. conversion to word \n 2. conversion to pdf \n 3. conversion to ppt"))



#CONVERSION TO WORD

if (ch==1):
     # create an instance of a word document 
     doc = docx.Document() 
  
     # add a heading of level 0 (largest heading) 
     doc.add_heading('Heading for the document', 0) 
  
     # add a paragraph and store  
     doc_para = doc.add_paragraph(p)
     # pictures can also be added to our word document 
     # width is optional 
     #doc.add_picture('test.jpg') 

     # now save the document to a location 
     doc.save('hello')




#CONVERSION TO PDF

if (ch==2):
    pdf = FPDF()
    pdf.add_page()
    with open("text.txt", 'rb') as fh:
            txt = fh.read().decode('latin-1')
        # Times 12
            pdf.set_font('Times', '', 20)
        # Output justified text
            pdf.multi_cell(0, 5, txt)
            pdf.ln(50)#gives line breaks
    pdf.output("multipage_simple.pdf")



# CONVERSION TO PPT

if (ch==3):
    
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Hello"
    subtitle.text = p

    prs.save('test.pptx')

    
    """
    f=open('temp.pptx')
    prs = Presentation(f)
    # Use the output from analyze_ppt to understand which layouts and placeholders
    # to use
    # Create a title slide first
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[3]
    title.text = "Quarterly Report"
    subtitle.text = "Generated on {:%m-%d-%Y}"
    subtitle.text = p
    prs.save('test.pptx')
    f.close()
    """


